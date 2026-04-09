"""PowerPoint-Dateien (.pptx) einlesen."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches


def parse_powerpoint(file_path: Path) -> list[dict]:
    """PowerPoint-Datei einlesen. Gibt Folien mit Text zurueck."""
    prs = Presentation(str(file_path))
    results = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        has_images = False
        has_charts = False
        has_tables = False

        for shape in slide.shapes:
            # Text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

            # Tabellen
            if shape.has_table:
                has_tables = True
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    header = " | ".join(rows[0])
                    body = "\n".join(" | ".join(r) for r in rows[1:])
                    texts.append(f"\n[Tabelle]\n{header}\n{'─' * len(header)}\n{body}")

            # Bilder
            if shape.shape_type == 13:  # Picture
                has_images = True

            # Charts
            if shape.has_chart:
                has_charts = True

        if texts:
            slide_text = "\n".join(texts)
            extras = []
            if has_images:
                extras.append("Bilder")
            if has_charts:
                extras.append("Diagramme")
            if has_tables:
                extras.append("Tabellen")

            results.append({
                "page": slide_num,
                "text": slide_text,
                "content_type": "text",
                "section": f"Folie {slide_num}",
                "contains": extras,
            })

    return results
